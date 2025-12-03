# scripts/generate_report.py

import os
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

# Optional AI integration (for openai>=1.0.0)
try:
    import openai
    AI_ENABLED = True
except ImportError:
    AI_ENABLED = False

# ---------------------------
# Paths / Folders
# ---------------------------
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_DIR = os.path.abspath(os.path.join(BASE_DIR, ".."))
DATA_DIR = os.path.join(PROJECT_DIR, "data")
REPORTS_DIR = os.path.join(PROJECT_DIR, "reports")

os.makedirs(REPORTS_DIR, exist_ok=True)


def _save_plot(fig, filename):
    """Helper to save matplotlib figure to reports folder."""
    out_path = os.path.join(REPORTS_DIR, filename)
    fig.tight_layout()
    fig.savefig(out_path)
    plt.close(fig)
    return out_path


def run_report(csv_path: str | None = None, generate_files: bool = True) -> dict:
    """
    Full pipeline:
    - Load CSV (or use data/current_dataset.csv or data/netflix_titles.csv)
    - Clean / engineer features
    - Generate plots
    - Create AI summary (if enabled)
    - Optionally generate PDF & PPT
    Returns dict with paths and summary.
    """

    # 1. Decide which CSV to use
    if csv_path is None:
        current_candidate = os.path.join(DATA_DIR, "current_dataset.csv")
        default_netflix = os.path.join(DATA_DIR, "netflix_titles.csv")

        if os.path.exists(current_candidate):
            csv_path = current_candidate
        else:
            csv_path = default_netflix

    if not os.path.exists(csv_path):
        raise FileNotFoundError(f"CSV file not found: {csv_path}")

    # 2. Load dataset
    df = pd.read_csv(csv_path)

    # 3. Cleaning / feature engineering
    df.columns = [c.strip() for c in df.columns]

    if "duration" in df.columns:
        df["duration_num"] = (
            df["duration"]
            .astype(str)
            .str.extract(r"(\d+)")
            .astype(float)
        )
    else:
        df["duration_num"] = None

    if "release_year" in df.columns:
        df["release_year"] = pd.to_numeric(df["release_year"], errors="coerce")

    sns.set(style="whitegrid")

    plot_files: dict[str, str] = {}

    # ---------------------------
    # Top genres
    # ---------------------------
    if "listed_in" in df.columns:
        top_genres = (
            df["listed_in"]
            .dropna()
            .astype(str)
            .str.split(", ", expand=True)
            .stack()
            .value_counts()
            .head(10)
        )

        if not top_genres.empty:
            fig, ax = plt.subplots(figsize=(8, 6))
            sns.barplot(x=top_genres.values, y=top_genres.index, dodge=False, ax=ax)
            ax.set_title("Top 10 Genres")
            ax.set_xlabel("Count")
            ax.set_ylabel("Genre")
            plot_files["top_genres"] = _save_plot(fig, "top_genres.png")
    else:
        top_genres = pd.Series(dtype=int)

    # ---------------------------
    # Top directors
    # ---------------------------
    if "director" in df.columns:
        top_directors = df["director"].dropna().value_counts().head(10)
        if not top_directors.empty:
            fig, ax = plt.subplots(figsize=(8, 6))
            sns.barplot(x=top_directors.values, y=top_directors.index, dodge=False, ax=ax)
            ax.set_title("Top 10 Directors")
            ax.set_xlabel("Count")
            ax.set_ylabel("Director")
            plot_files["top_directors"] = _save_plot(fig, "top_directors.png")
    else:
        top_directors = pd.Series(dtype=int)

    # ---------------------------
    # Ratings distribution
    # ---------------------------
    if "rating" in df.columns:
        rating_counts = df["rating"].value_counts()
        if not rating_counts.empty:
            fig, ax = plt.subplots(figsize=(6, 5))
            sns.countplot(
                y="rating",
                data=df,
                order=rating_counts.index,
                dodge=False,
                ax=ax,
            )
            ax.set_title("Rating Distribution")
            ax.set_xlabel("Count")
            ax.set_ylabel("Rating")
            plot_files["rating_distribution"] = _save_plot(fig, "rating_distribution.png")
    else:
        rating_counts = pd.Series(dtype=int)

    # ---------------------------
    # Titles per year (trend)
    # ---------------------------
    if "release_year" in df.columns:
        per_year = df["release_year"].value_counts().sort_index()
        if not per_year.empty:
            fig, ax = plt.subplots(figsize=(8, 5))
            per_year.plot(kind="line", marker="o", ax=ax)
            ax.set_title("Titles Released per Year")
            ax.set_xlabel("Year")
            ax.set_ylabel("Number of Titles")
            plot_files["titles_per_year"] = _save_plot(fig, "titles_per_year.png")

    # ---------------------------
    # Type distribution (Movie vs TV Show)
    # ---------------------------
    if "type" in df.columns:
        type_counts = df["type"].value_counts()
        if not type_counts.empty:
            fig, ax = plt.subplots(figsize=(6, 5))
            sns.barplot(x=type_counts.index, y=type_counts.values, dodge=False, ax=ax)
            ax.set_title("Title Type Distribution")
            ax.set_xlabel("Type")
            ax.set_ylabel("Count")
            plot_files["type_count"] = _save_plot(fig, "type_count.png")

    # ---------------------------
    # Average duration by type (if available)
    # ---------------------------
    if "duration_num" in df.columns and "type" in df.columns:
        avg_duration = df.dropna(subset=["duration_num"]).groupby("type")["duration_num"].mean()
        if not avg_duration.empty:
            fig, ax = plt.subplots(figsize=(6, 5))
            sns.barplot(x=avg_duration.index, y=avg_duration.values, dodge=False, ax=ax)
            ax.set_title("Average Duration by Type")
            ax.set_xlabel("Type")
            ax.set_ylabel("Average Duration")
            plot_files["avg_movie_duration"] = _save_plot(fig, "avg_movie_duration.png")

    # ---------------------------
    # AI Executive Summary
    # ---------------------------
    ai_summary = ""

    if AI_ENABLED:
        try:
            summary_prompt = f"""
            You are a data analyst. Summarize this Netflix dataset analysis.

            - Top genres (with counts): {top_genres.to_dict()}
            - Top directors (with counts): {top_directors.to_dict()}
            - Ratings distribution: {rating_counts.to_dict()}

            Write a short executive summary in 4–5 sentences.
            Focus on:
            - Which genres dominate
            - What the rating distribution suggests about target audience
            - Any diversity in directors or content
            Avoid technical jargon. Write clearly for business stakeholders.
            """

            response = openai.chat.completions.create(
                model="gpt-4o-mini",
                messages=[{"role": "user", "content": summary_prompt}],
                max_tokens=250,
            )
            ai_summary = response.choices[0].message.content.strip()
            print("✅ AI Summary generated.")
        except Exception as e:
            print("⚠️ AI summary not generated:", e)
    else:
        print("AI integration not enabled. Install 'openai' and set OPENAI_API_KEY to use this feature.")

    pdf_file = None
    ppt_file = None

    if generate_files:
        # ---------------------------
        # Generate PDF report
        # ---------------------------
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", "B", 16)
        pdf.cell(0, 10, "Netflix Analysis Report", ln=True, align="C")
        pdf.ln(10)

        if ai_summary:
            pdf.set_font("Arial", "", 12)
            pdf.multi_cell(0, 8, ai_summary)
            pdf.ln(10)

        for key in [
            "top_genres",
            "top_directors",
            "rating_distribution",
            "titles_per_year",
            "type_count",
            "avg_movie_duration",
        ]:
            if key in plot_files:
                pdf.image(plot_files[key], w=180)
                pdf.ln(8)

        pdf_file = os.path.join(REPORTS_DIR, "netflix_report.pdf")
        pdf.output(pdf_file)
        print(f"✅ PDF report generated: {pdf_file}")

        # ---------------------------
        # Generate PowerPoint report
        # ---------------------------
        ppt = Presentation()
        blank_layout = ppt.slide_layouts[5]  # blank

        slide = ppt.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.3), Inches(8.0), Inches(1.0))
        title_frame = title_box.text_frame
        title_frame.text = "Executive Summary"
        title_frame.paragraphs[0].font.size = Pt(28)

        if ai_summary:
            body_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.0), Inches(4.5))
            body_frame = body_box.text_frame
            body_frame.text = ai_summary
            body_frame.paragraphs[0].font.size = Pt(16)

        for key in [
            "top_genres",
            "top_directors",
            "rating_distribution",
            "titles_per_year",
            "type_count",
            "avg_movie_duration",
        ]:
            if key in plot_files:
                slide = ppt.slides.add_slide(blank_layout)
                slide.shapes.add_picture(plot_files[key], Inches(0.7), Inches(0.7), width=Inches(8.5))

        ppt_file = os.path.join(REPORTS_DIR, "netflix_report.pptx")
        ppt.save(ppt_file)
        print(f"✅ PowerPoint report generated: {ppt_file}")

        print("🎉 All tasks completed! Check 'reports/' for plots, PDF, PPT, and AI summary.")

    return {
        "csv_path": csv_path,
        "pdf_file": pdf_file,
        "ppt_file": ppt_file,
        "ai_summary": ai_summary,
        "plots": plot_files,
    }


if __name__ == "__main__":
    # Run locally from terminal:
    # (venv) python scripts/generate_report.py
    result = run_report()
    print("\nAI Summary:\n", result["ai_summary"])
